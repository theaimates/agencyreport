"""
ReportFlow — Pitch Deck Generator
Outputs: .tmp/ReportFlow_Pitch_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUT_DIR = os.path.join(os.path.dirname(__file__), '..', '.tmp')
os.makedirs(OUT_DIR, exist_ok=True)
OUT_PATH = os.path.join(OUT_DIR, 'ReportFlow_Pitch_Deck.pptx')

# ── Colour palette (matches dashboard) ──────────────────────────
BG       = RGBColor(0x0E, 0x0D, 0x0B)   # near-black
ACCENT   = RGBColor(0xD4, 0xA8, 0x43)   # gold
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x9A, 0x96, 0x90)
INDIGO   = RGBColor(0x63, 0x66, 0xF1)
GREEN    = RGBColor(0x3E, 0xCF, 0x8E)
CARD_BG  = RGBColor(0x1A, 0x18, 0x15)
RED      = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # completely blank


# ── Helpers ──────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def text_box(slide, text, x, y, w, h,
             font_size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def rect(slide, x, y, w, h, fill=CARD_BG, line=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def accent_line(slide, x, y, w):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def multiline_box(slide, lines, x, y, w, h, font_size=18, color=WHITE, font_name='Calibri'):
    """lines: list of (text, bold, color_override)"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col or color
        run.font.name = font_name


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
# Gold top bar
rect(s, 0, 0, 13.33, 0.08, fill=ACCENT)
# Logo mark
rect(s, 0.7, 2.8, 0.65, 0.65, fill=ACCENT)
text_box(s, 'R', 0.7, 2.78, 0.65, 0.65, font_size=28, bold=True, color=BG, align=PP_ALIGN.CENTER)
text_box(s, 'ReportFlow', 1.5, 2.75, 5, 0.7, font_size=40, bold=True, color=WHITE)
text_box(s, 'Automated AI Reporting for Marketing Agencies', 1.5, 3.5, 9, 0.6, font_size=20, color=ACCENT)
text_box(s, 'Stop building reports. Start closing clients.', 1.5, 4.3, 9, 0.5, font_size=16, color=MUTED, italic=True)
text_box(s, '2026', 1.5, 6.8, 2, 0.4, font_size=13, color=MUTED)
accent_line(s, 1.5, 4.15, 4)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════════════════════════
s = add_slide()
accent_line(s, 0.7, 1.1, 1.2)
text_box(s, 'The Problem', 0.7, 1.2, 8, 0.7, font_size=36, bold=True, color=WHITE)
text_box(s, 'Every month, agency account managers lose days to manual reporting.', 0.7, 2.1, 10, 0.6, font_size=18, color=MUTED)

problems = [
    ('⏱  4–6 hours per client, per month', ACCENT),
    ('📊  Copy-pasting from GA4, Meta, Google Ads, Stripe into Google Docs', WHITE),
    ('🔁  Same report, different numbers, every single month', WHITE),
    ('😤  Clients get PDFs that are already out of date by send time', WHITE),
    ('💸  At $150/hr, that\'s $600–900 of billable time thrown away per client', RED),
]
for i, (txt, col) in enumerate(problems):
    rect(s, 0.7, 2.9 + i * 0.82, 11.5, 0.68, fill=CARD_BG, line=RGBColor(0x2A, 0x28, 0x24))
    text_box(s, txt, 1.0, 2.95 + i * 0.82, 11, 0.55, font_size=16, color=col)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — The Old Way vs ReportFlow
# ════════════════════════════════════════════════════════════════
s = add_slide()
accent_line(s, 0.7, 0.6, 1.2)
text_box(s, 'The Old Way vs. ReportFlow', 0.7, 0.7, 11, 0.7, font_size=32, bold=True, color=WHITE)

# Left column — Old Way
rect(s, 0.5, 1.6, 5.8, 5.3, fill=RGBColor(0x18, 0x10, 0x10), line=RED)
text_box(s, '❌  The Old Way', 0.9, 1.75, 5, 0.5, font_size=18, bold=True, color=RED)
old = [
    'Pull data from 5 different platforms',
    'Manually calculate changes & trends',
    'Write narrative summaries from scratch',
    'Format everything in Google Slides/Docs',
    'Send and pray the client reads it',
    'Do it all again next month',
]
for i, t in enumerate(old):
    text_box(s, f'  •  {t}', 0.8, 2.35 + i * 0.65, 5.3, 0.55, font_size=13, color=MUTED)

# Right column — ReportFlow
rect(s, 6.9, 1.6, 5.8, 5.3, fill=RGBColor(0x10, 0x18, 0x12), line=GREEN)
text_box(s, '✅  ReportFlow', 7.3, 1.75, 5, 0.5, font_size=18, bold=True, color=GREEN)
new = [
    'Data pulled automatically every day',
    'Trends & changes calculated instantly',
    'GPT-4o writes the narrative for you',
    'Live dashboard, always up to date',
    'Email sent automatically every Monday',
    'Zero manual work after setup',
]
for i, t in enumerate(new):
    text_box(s, f'  •  {t}', 7.2, 2.35 + i * 0.65, 5.3, 0.55, font_size=13, color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — How It Works
# ════════════════════════════════════════════════════════════════
s = add_slide()
accent_line(s, 0.7, 0.6, 1.2)
text_box(s, 'How It Works', 0.7, 0.7, 8, 0.7, font_size=32, bold=True, color=WHITE)
text_box(s, 'Five automated workflows. Zero manual steps.', 0.7, 1.5, 10, 0.5, font_size=16, color=MUTED)

steps = [
    ('1', 'Collect',    'Every morning at 6 AM,\nReportFlow pulls from GA4, Meta,\nGoogle Ads, Stripe & Mailchimp', INDIGO),
    ('2', 'Store',      'All data lands in Airtable —\nclean, structured, and ready\nfor analysis', ACCENT),
    ('3', 'Analyse',    'GPT-4o reads the week\'s data\nand writes a personalised\nmarketing analysis', GREEN),
    ('4', 'Deliver',    'A branded email report is\nautomatically sent to the client\nevery Monday morning', ACCENT),
    ('5', 'Dashboard',  'The client (or you) can log in\nany time to see live charts\nand drill into campaigns', INDIGO),
]

for i, (num, title, body, col) in enumerate(steps):
    x = 0.5 + i * 2.55
    rect(s, x, 2.1, 2.35, 4.8, fill=CARD_BG, line=col)
    # Number circle
    rect(s, x + 0.85, 2.25, 0.65, 0.58, fill=col)
    text_box(s, num, x + 0.85, 2.22, 0.65, 0.58, font_size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
    text_box(s, title, x + 0.1, 3.0, 2.15, 0.5, font_size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
    text_box(s, body, x + 0.1, 3.55, 2.15, 2.8, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Why Not Just Use ChatGPT?
# ════════════════════════════════════════════════════════════════
s = add_slide()
accent_line(s, 0.7, 0.6, 1.2)
text_box(s, '"Why not just use ChatGPT?"', 0.7, 0.7, 11, 0.7, font_size=32, bold=True, color=WHITE)
text_box(s, 'Great question. Here\'s the difference.', 0.7, 1.5, 8, 0.5, font_size=16, color=MUTED)

comparisons = [
    ('ChatGPT / Claude', 'ReportFlow'),
    ('You manually copy-paste data into the chat', 'Data is pulled automatically every day'),
    ('You must remember to ask for the report', 'Reports are delivered on a schedule — no prompting'),
    ('No memory between sessions', 'Full historical data in Airtable — trends over months'),
    ('You format and send the output yourself', 'Branded email delivered directly to the client'),
    ('No live dashboard', 'Always-on dashboard the client can log into'),
    ('One-off output', 'Continuous, compounding intelligence over time'),
]

rect(s, 0.5, 2.1, 5.9, 0.5, fill=RGBColor(0x1E, 0x1C, 0x18))
rect(s, 6.9, 2.1, 5.9, 0.5, fill=RGBColor(0x10, 0x18, 0x12))
text_box(s, comparisons[0][0], 0.7, 2.15, 5.5, 0.4, font_size=14, bold=True, color=MUTED)
text_box(s, comparisons[0][1], 7.1, 2.15, 5.5, 0.4, font_size=14, bold=True, color=GREEN)

for i, (left, right) in enumerate(comparisons[1:]):
    y = 2.75 + i * 0.65
    text_box(s, f'✗  {left}', 0.7, y, 5.8, 0.55, font_size=12, color=MUTED)
    text_box(s, f'✓  {right}', 7.1, y, 5.5, 0.55, font_size=12, color=WHITE)
    # divider
    rect(s, 0.5, y + 0.58, 12.3, 0.02, fill=RGBColor(0x2A, 0x28, 0x24))


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — The Dashboard
# ════════════════════════════════════════════════════════════════
s = add_slide()
accent_line(s, 0.7, 0.6, 1.2)
text_box(s, 'A Dashboard Your Clients Actually Want to Open', 0.7, 0.7, 11, 0.7, font_size=28, bold=True, color=WHITE)

features = [
    ('📈', 'Revenue Trend',      'Weekly chart with period-over-period comparison'),
    ('🍩', 'Channel Mix',        'Donut chart showing revenue split by platform'),
    ('🏆', 'Channel Performance','ROAS, spend, revenue, and conversion rate per channel'),
    ('📋', 'Campaign Table',     'Every campaign with CTR, CPC, impressions, and status'),
    ('🤖', 'AI Analysis',        'GPT-4o written summary: biggest win, watch closely, next move'),
    ('📅', 'Date Ranges',        'This month, last month, Q1, last 90 days, YTD'),
]
for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.25
    y = 1.8 + row * 2.4
    rect(s, x, y, 3.9, 2.1, fill=CARD_BG, line=RGBColor(0x2A, 0x28, 0x24))
    text_box(s, icon, x + 0.2, y + 0.2, 0.6, 0.6, font_size=24)
    text_box(s, title, x + 0.9, y + 0.2, 2.8, 0.5, font_size=14, bold=True, color=ACCENT)
    text_box(s, desc, x + 0.2, y + 0.85, 3.5, 1.0, font_size=12, color=MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — ROI / Value Proof
# ════════════════════════════════════════════════════════════════
s = add_slide()
accent_line(s, 0.7, 0.6, 1.2)
text_box(s, 'The Numbers Make the Decision Easy', 0.7, 0.7, 11, 0.7, font_size=32, bold=True, color=WHITE)

# Big stats
stats = [
    ('4–6 hrs',  'saved per client\nper month'),
    ('×10',      'faster reporting\nthan manual'),
    ('100%',     'on-time delivery\nevery week'),
    ('$0',       'extra headcount\nneeded'),
]
for i, (big, small) in enumerate(stats):
    x = 0.5 + i * 3.1
    rect(s, x, 1.7, 2.8, 2.0, fill=CARD_BG, line=ACCENT)
    text_box(s, big, x + 0.1, 1.8, 2.6, 0.9, font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    text_box(s, small, x + 0.1, 2.7, 2.6, 0.8, font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

text_box(s, 'Real scenario: 10 clients × 5 hrs/month = 50 hrs of manual work', 0.7, 4.0, 11, 0.5, font_size=15, color=MUTED)
text_box(s, 'With ReportFlow: those 50 hours become 0.', 0.7, 4.55, 11, 0.5, font_size=18, bold=True, color=WHITE)
text_box(s, 'At $150/hr that\'s $7,500/month in recovered billable capacity — per agency.', 0.7, 5.15, 11, 0.5, font_size=15, color=ACCENT)

rect(s, 0.7, 5.85, 11.5, 0.95, fill=RGBColor(0x10, 0x18, 0x12), line=GREEN)
text_box(s, '✓  That\'s not a productivity improvement. That\'s a business model change.', 1.0, 5.98, 11, 0.6, font_size=15, bold=True, color=GREEN)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Who It's For
# ════════════════════════════════════════════════════════════════
s = add_slide()
accent_line(s, 0.7, 0.6, 1.2)
text_box(s, 'Built For', 0.7, 0.7, 8, 0.7, font_size=36, bold=True, color=WHITE)

personas = [
    ('🏢', 'Marketing Agencies',
     '5–50 clients\nRunning paid ads across Meta, Google, and email\nTired of the monthly reporting grind'),
    ('👤', 'Solo Consultants',
     'Managing 3–15 clients alone\nNeed to look like a full agency\nWant to impress clients without the overhead'),
    ('📈', 'Growth Teams',
     'In-house marketing teams\nReporting to executives weekly\nNeed clean data without a BI tool budget'),
]
for i, (icon, title, body) in enumerate(personas):
    x = 0.5 + i * 4.2
    rect(s, x, 1.7, 3.9, 5.2, fill=CARD_BG, line=ACCENT)
    text_box(s, icon, x + 1.5, 1.9, 0.9, 0.9, font_size=36, align=PP_ALIGN.CENTER)
    text_box(s, title, x + 0.2, 2.9, 3.5, 0.55, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    text_box(s, body, x + 0.3, 3.55, 3.3, 2.8, font_size=13, color=MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Tech Stack (Credibility)
# ════════════════════════════════════════════════════════════════
s = add_slide()
accent_line(s, 0.7, 0.6, 1.2)
text_box(s, 'Built on Proven Infrastructure', 0.7, 0.7, 10, 0.7, font_size=32, bold=True, color=WHITE)
text_box(s, 'No custom servers to maintain. No vendor lock-in. All tools you already trust.', 0.7, 1.55, 11, 0.5, font_size=15, color=MUTED)

stack = [
    ('n8n',        'Workflow Automation',  'Self-hosted or cloud. 6 pre-built workflows handle all data collection, aggregation, AI generation, and delivery.', INDIGO),
    ('Airtable',   'Data Layer',           'Structured storage for clients, KPIs, campaigns, revenue, channel performance, and AI summaries.', GREEN),
    ('GPT-4o',     'AI Engine',            'OpenAI\'s most capable model writes personalised weekly analysis for each client automatically.', ACCENT),
    ('HTML/JS',    'Dashboard',            'Lightweight, dependency-free frontend. Deployable anywhere. No React, no build step, no maintenance.', MUTED),
]
for i, (name, role, desc, col) in enumerate(stack):
    y = 2.3 + i * 1.15
    rect(s, 0.5, y, 12.2, 0.95, fill=CARD_BG, line=col)
    text_box(s, name, 0.75, y + 0.08, 1.6, 0.6, font_size=16, bold=True, color=col)
    text_box(s, role, 2.5, y + 0.08, 2.8, 0.4, font_size=12, bold=True, color=WHITE)
    text_box(s, desc, 5.5, y + 0.08, 7.0, 0.7, font_size=11, color=MUTED)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — What You Get
# ════════════════════════════════════════════════════════════════
s = add_slide()
accent_line(s, 0.7, 0.6, 1.2)
text_box(s, 'Everything Included', 0.7, 0.7, 10, 0.7, font_size=36, bold=True, color=WHITE)

deliverables = [
    '6 production-ready n8n workflow JSON files (import in 2 minutes)',
    'Pre-built Airtable schema — Clients, KPIs, Campaigns, Revenue, Channel Performance, AI Summaries',
    'Dashboard frontend — deployable to any static host (Netlify, Vercel, GitHub Pages)',
    'Full setup guide — step-by-step from zero to first automated report',
    'Webhook API layer — secure, paginated, with CORS protection',
    'Mock data mode — demostrate to clients before live APIs are connected',
]
for i, item in enumerate(deliverables):
    rect(s, 0.5, 1.65 + i * 0.82, 12.2, 0.68, fill=CARD_BG, line=GREEN)
    text_box(s, f'✓  {item}', 0.8, 1.72 + i * 0.82, 11.5, 0.55, font_size=14, color=WHITE)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Call to Action
# ════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=BG)
rect(s, 0, 0, 13.33, 0.08, fill=ACCENT)
rect(s, 0, 7.42, 13.33, 0.08, fill=ACCENT)

text_box(s, 'Ready to automate your reporting?', 1.5, 1.5, 10.5, 0.9, font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
accent_line(s, 4.5, 2.6, 4.3)
text_box(s, 'Set up once. Report forever.', 1.5, 2.8, 10.5, 0.7, font_size=22, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
text_box(s, 'ReportFlow is ready to deploy today.\nNo custom development. No subscriptions. Your data, your infrastructure.', 1.5, 3.7, 10.5, 1.0, font_size=16, color=MUTED, align=PP_ALIGN.CENTER)

# CTA box
rect(s, 3.5, 5.0, 6.3, 1.0, fill=ACCENT)
text_box(s, 'Book a Setup Call  →', 3.5, 5.05, 6.3, 0.85, font_size=22, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f'Saved: {OUT_PATH}')
